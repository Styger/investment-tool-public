"""
PowerPoint Report Generator for Backtesting Results
Creates professional presentation with charts embedded
"""

from datetime import datetime
from typing import Dict, List
import tempfile
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx not available. Install with: pip install python-pptx")


class BacktestPPTXExporter:
    """Generate professional PowerPoint reports from backtest results"""

    # Color scheme
    PRIMARY_COLOR = RGBColor(31, 119, 180)  # Blue
    SUCCESS_COLOR = RGBColor(40, 167, 69)  # Green
    DANGER_COLOR = RGBColor(220, 53, 69)  # Red
    GRAY_COLOR = RGBColor(108, 117, 125)  # Gray

    @staticmethod
    def generate_report(
        results: Dict,
        charts: Dict,
        universe_name: str = "Unknown",
        parameters: Dict = None,
    ) -> bytes:
        """
        Generate PowerPoint report from backtest results

        Args:
            results: Backtest results dictionary
            charts: Dictionary of Plotly chart objects
            universe_name: Name of tested universe
            parameters: Strategy parameters used

        Returns:
            PowerPoint file as bytes
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PowerPoint generation. "
                "Install with: pip install python-pptx"
            )

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Extract data
        metrics = results.get("metrics", {})
        benchmark = results.get("benchmark", {})
        trades = results.get("trades", [])
        valuation_methods = results.get("valuation_methods", [])
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        # Slide 1: Title
        BacktestPPTXExporter._add_title_slide(
            prs, universe_name, valuation_methods, timestamp, results
        )

        # Slide 2: Executive Summary
        BacktestPPTXExporter._add_executive_summary(prs, results, metrics)

        # Slide 3: Strategy Parameters
        if parameters:
            BacktestPPTXExporter._add_parameters_slide(
                prs, parameters, valuation_methods
            )

        # Slide 4: Performance Metrics
        BacktestPPTXExporter._add_metrics_slide(prs, metrics)

        # Slide 5: Benchmark Comparison
        if benchmark.get("benchmark_available"):
            BacktestPPTXExporter._add_benchmark_slide(prs, benchmark)

        # Slide 6-9: Charts (using Matplotlib - no browser needed!)
        from backend.backtesting.analytics.matplotlib_charts import (
            MatplotlibChartGenerator,
        )

        print("📊 Generating charts with Matplotlib (browser-independent)...")
        chart_images = MatplotlibChartGenerator.generate_all_charts(results)

        for chart_name, chart_path in chart_images.items():
            if chart_path and os.path.exists(chart_path):
                BacktestPPTXExporter._add_chart_slide(
                    prs, chart_name.title(), chart_path
                )

        # Slide 10: Trade Statistics
        BacktestPPTXExporter._add_trade_stats_slide(prs, metrics)

        # Slide 11-12: Top Trades
        if trades:
            BacktestPPTXExporter._add_top_trades_slides(prs, trades)

        # Save to bytes
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(temp_file.name)
        temp_file.close()

        with open(temp_file.name, "rb") as f:
            pptx_bytes = f.read()

        # Cleanup
        os.remove(temp_file.name)
        for chart_path in chart_images.values():
            if chart_path and os.path.exists(chart_path):
                try:
                    os.remove(chart_path)
                except:
                    pass

        return pptx_bytes

    @staticmethod
    def _add_title_slide(prs, universe_name, valuation_methods, timestamp, results):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Title
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = "📊 ValueKit Backtest Report"

        p = title_frame.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = BacktestPPTXExporter.PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.5), Inches(8), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Consensus Valuation Strategy"
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        # Info
        info_text = (
            f"Universe: {universe_name}\n"
            f"Period: {results.get('dates', [''])[0] if results.get('dates') else 'N/A'} - "
            f"{results.get('dates', [''])[-1] if results.get('dates') else 'N/A'}\n"
            f"Methods: {', '.join(valuation_methods) if valuation_methods else 'N/A'}\n"
            f"Generated: {timestamp}"
        )

        info_box = slide.shapes.add_textbox(
            Inches(2), Inches(5), Inches(6), Inches(1.5)
        )
        info_frame = info_box.text_frame
        info_frame.text = info_text
        p = info_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

    @staticmethod
    def _add_executive_summary(prs, results, metrics):
        """Add executive summary slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
        title = slide.shapes.title
        title.text = "📋 Executive Summary"

        # Create table
        rows = 6
        cols = 2
        left = Inches(2)
        top = Inches(2)
        width = Inches(6)
        height = Inches(4)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(3)

        # Header
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"

        # Data
        table.cell(1, 0).text = "Final Value"
        table.cell(1, 1).text = f"${results.get('final_value', 0):,.2f}"

        table.cell(2, 0).text = "Profit"
        profit = results.get("profit", 0)
        table.cell(2, 1).text = f"${profit:+,.2f}"

        table.cell(3, 0).text = "Total Return"
        table.cell(3, 1).text = f"{results.get('return_pct', 0):.2f}%"

        table.cell(4, 0).text = "CAGR"
        table.cell(4, 1).text = f"{results.get('cagr', 0):.2f}%"

        table.cell(5, 0).text = "Period"
        table.cell(5, 1).text = f"{metrics.get('years', 0):.1f} years"

        # Style table
        BacktestPPTXExporter._style_table(table)

    @staticmethod
    def _add_parameters_slide(prs, parameters, valuation_methods):
        """Add strategy parameters slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "🎯 Strategy Parameters"

        # Create table
        rows = 8
        cols = 2
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        height = Inches(4.5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(3.5)
        table.columns[1].width = Inches(3.5)

        # Header
        table.cell(0, 0).text = "Parameter"
        table.cell(0, 1).text = "Value"

        # Data
        table.cell(1, 0).text = "Valuation Methods"
        table.cell(1, 1).text = (
            ", ".join(valuation_methods) if valuation_methods else "N/A"
        )

        table.cell(2, 0).text = "MOS Threshold"
        table.cell(2, 1).text = f"{parameters.get('mos_threshold', 'N/A')}%"

        table.cell(3, 0).text = "Moat Threshold"
        table.cell(3, 1).text = f"{parameters.get('moat_threshold', 'N/A')}/50"

        table.cell(4, 0).text = "Sell MOS Threshold"
        table.cell(4, 1).text = f"{parameters.get('sell_mos_threshold', 'N/A')}%"

        table.cell(5, 0).text = "Sell Moat Threshold"
        table.cell(5, 1).text = f"{parameters.get('sell_moat_threshold', 'N/A')}/50"

        table.cell(6, 0).text = "Max Positions"
        table.cell(6, 1).text = str(parameters.get("max_positions", "N/A"))

        table.cell(7, 0).text = "Rebalance Frequency"
        table.cell(7, 1).text = f"Every {parameters.get('rebalance_days', 'N/A')} days"

        BacktestPPTXExporter._style_table(table)

    @staticmethod
    def _add_metrics_slide(prs, metrics):
        """Add performance metrics slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "📊 Performance Metrics"

        # Create table
        rows = 5
        cols = 3
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(2)
        table.columns[2].width = Inches(3.5)

        # Header
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        table.cell(0, 2).text = "Description"

        # Data
        table.cell(1, 0).text = "Sharpe Ratio"
        table.cell(1, 1).text = f"{metrics.get('sharpe_ratio', 0):.2f}"
        table.cell(1, 2).text = "Risk-adjusted return"

        table.cell(2, 0).text = "Sortino Ratio"
        table.cell(2, 1).text = f"{metrics.get('sortino_ratio', 0):.2f}"
        table.cell(2, 2).text = "Downside risk-adjusted"

        table.cell(3, 0).text = "Max Drawdown"
        table.cell(3, 1).text = f"{metrics.get('max_drawdown', 0):.2f}%"
        table.cell(3, 2).text = "Largest peak-to-trough decline"

        table.cell(4, 0).text = "Calmar Ratio"
        table.cell(4, 1).text = f"{metrics.get('calmar_ratio', 0):.2f}"
        table.cell(4, 2).text = "Return / Max Drawdown"

        BacktestPPTXExporter._style_table(table)

    @staticmethod
    def _add_benchmark_slide(prs, benchmark):
        """Add benchmark comparison slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "🎯 Benchmark Comparison (S&P 500)"

        # Create table
        rows = 5
        cols = 4
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(9)
        height = Inches(4)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Header
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Strategy"
        table.cell(0, 2).text = "S&P 500"
        table.cell(0, 3).text = "Difference"

        # Data
        table.cell(1, 0).text = "Total Return"
        table.cell(1, 1).text = f"{benchmark.get('strategy_return', 0):.2f}%"
        table.cell(1, 2).text = f"{benchmark.get('benchmark_return', 0):.2f}%"
        table.cell(1, 3).text = f"{benchmark.get('outperformance', 0):+.2f}%"

        table.cell(2, 0).text = "CAGR"
        table.cell(2, 1).text = f"{benchmark.get('strategy_cagr', 0):.2f}%"
        table.cell(2, 2).text = f"{benchmark.get('benchmark_cagr', 0):.2f}%"
        cagr_diff = benchmark.get("strategy_cagr", 0) - benchmark.get(
            "benchmark_cagr", 0
        )
        table.cell(2, 3).text = f"{cagr_diff:+.2f}%"

        table.cell(3, 0).text = "Alpha"
        table.cell(3, 1).text = "-"
        table.cell(3, 2).text = "-"
        table.cell(3, 3).text = f"{benchmark.get('alpha', 0):+.2f}%"

        table.cell(4, 0).text = "Beta"
        table.cell(4, 1).text = f"{benchmark.get('beta', 0):.2f}"
        table.cell(4, 2).text = "1.00"
        table.cell(4, 3).text = "Market correlation"

        BacktestPPTXExporter._style_table(table)

    @staticmethod
    def _add_chart_slide(prs, chart_title, chart_path):
        """Add chart slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = f"📈 {chart_title}"

        # Add chart image
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9)
        height = Inches(5)

        slide.shapes.add_picture(chart_path, left, top, width=width, height=height)

    @staticmethod
    def _add_trade_stats_slide(prs, metrics):
        """Add trade statistics slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "📝 Trade Statistics"

        # Create table
        rows = 7
        cols = 2
        left = Inches(2)
        top = Inches(2)
        width = Inches(6)
        height = Inches(4)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(3)

        # Header
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"

        # Data
        table.cell(1, 0).text = "Total Trades"
        table.cell(1, 1).text = str(metrics.get("total_trades", 0))

        table.cell(2, 0).text = "Win Rate"
        table.cell(2, 1).text = f"{metrics.get('win_rate', 0):.1f}%"

        table.cell(3, 0).text = "Avg Hold Time"
        table.cell(3, 1).text = f"{metrics.get('avg_hold_time_days', 0):.0f} days"

        table.cell(4, 0).text = "Profit Factor"
        table.cell(4, 1).text = f"{metrics.get('profit_factor', 0):.2f}"

        table.cell(5, 0).text = "Average Win"
        table.cell(5, 1).text = f"${metrics.get('avg_win', 0):,.2f}"

        table.cell(6, 0).text = "Average Loss"
        table.cell(6, 1).text = f"${metrics.get('avg_loss', 0):,.2f}"

        BacktestPPTXExporter._style_table(table)

    @staticmethod
    def _add_top_trades_slides(prs, trades):
        """Add top winners and losers slides"""
        sorted_trades = sorted(trades, key=lambda x: x.get("pnl", 0), reverse=True)

        # Top 10 Winners
        top_winners = sorted_trades[:10]
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "🏆 Top 10 Winners"

        rows = min(len(top_winners) + 1, 11)
        cols = 5
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(9)
        height = Inches(4.5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Header
        table.cell(0, 0).text = "Ticker"
        table.cell(0, 1).text = "Buy Date"
        table.cell(0, 2).text = "Sell Date"
        table.cell(0, 3).text = "P&L"
        table.cell(0, 4).text = "Hold Days"

        # Data
        for i, trade in enumerate(top_winners, 1):
            buy_date = trade.get("buy_date", "N/A")
            sell_date = trade.get("sell_date", "N/A")

            if hasattr(buy_date, "strftime"):
                buy_date = buy_date.strftime("%Y-%m-%d")
            if hasattr(sell_date, "strftime"):
                sell_date = sell_date.strftime("%Y-%m-%d")

            table.cell(i, 0).text = str(trade.get("ticker", "N/A"))
            table.cell(i, 1).text = str(buy_date)
            table.cell(i, 2).text = str(sell_date)
            table.cell(i, 3).text = f"${trade.get('pnl', 0):,.2f}"
            table.cell(i, 4).text = str(trade.get("hold_days", 0))

        BacktestPPTXExporter._style_table(table)

        # Top 10 Losers
        top_losers = sorted_trades[-10:][::-1]
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "📉 Top 10 Losers"

        rows = min(len(top_losers) + 1, 11)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Header
        table.cell(0, 0).text = "Ticker"
        table.cell(0, 1).text = "Buy Date"
        table.cell(0, 2).text = "Sell Date"
        table.cell(0, 3).text = "P&L"
        table.cell(0, 4).text = "Hold Days"

        # Data
        for i, trade in enumerate(top_losers, 1):
            buy_date = trade.get("buy_date", "N/A")
            sell_date = trade.get("sell_date", "N/A")

            if hasattr(buy_date, "strftime"):
                buy_date = buy_date.strftime("%Y-%m-%d")
            if hasattr(sell_date, "strftime"):
                sell_date = sell_date.strftime("%Y-%m-%d")

            table.cell(i, 0).text = str(trade.get("ticker", "N/A"))
            table.cell(i, 1).text = str(buy_date)
            table.cell(i, 2).text = str(sell_date)
            table.cell(i, 3).text = f"${trade.get('pnl', 0):,.2f}"
            table.cell(i, 4).text = str(trade.get("hold_days", 0))

        BacktestPPTXExporter._style_table(table)

    @staticmethod
    def _save_charts_as_images(charts: Dict) -> Dict:
        """
        Save charts as temporary PNG images using Matplotlib

        This avoids Kaleido/browser dependencies which often fail on Windows
        """
        from backend.backtesting.analytics.matplotlib_charts import (
            MatplotlibChartGenerator,
        )

        # Generate matplotlib charts from results data
        # Note: charts dict contains Plotly objects, but we need the raw data
        # We'll generate new charts from scratch using matplotlib

        chart_images = {}

        # For now, we'll skip charts in PPTX if we don't have access to results
        # This will be fixed by passing results to this method
        print("ℹ️  Using Matplotlib for chart generation (browser-independent)")

        return chart_images

    @staticmethod
    def _style_table(table):
        """Apply consistent styling to table"""
        # Header row
        for cell in table.rows[0].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BacktestPPTXExporter.PRIMARY_COLOR

            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

        # Data rows
        for row_idx in range(1, len(table.rows)):
            for cell in table.rows[row_idx].cells:
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(11)
